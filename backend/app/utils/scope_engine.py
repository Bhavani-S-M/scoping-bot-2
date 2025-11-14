# app/utils/scope_engine.py
from __future__ import annotations
import json, re, logging, math, os, tempfile,anyio,pytesseract, openpyxl,tiktoken, pytz, graphviz,requests
from app import models
from calendar import monthrange
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation
from io import BytesIO
from PIL import Image
from app.config.config import QDRANT_COLLECTION
from typing import Dict, Any, List
from datetime import datetime, timedelta
from app.utils import azure_blob
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from sqlalchemy.orm import selectinload
from app.utils.ai_clients import (
    get_llm_client,
    get_embed_client,
    get_qdrant_client,
    embed_text_ollama,
)


logger = logging.getLogger(__name__)

# Init AI services
llm_cfg = get_llm_client()
embed_cfg = get_embed_client()
qdrant = get_qdrant_client()

def ollama_chat(prompt: str, model: str = llm_cfg["model"], temperature: float = 0.7) -> str:
    """Call Ollama to generate text from a prompt."""
    try:
        resp = requests.post(
            f"{llm_cfg['host']}/api/generate",
            json={"model": model, "prompt": prompt, "temperature": temperature, "stream": False},
        )
        resp.raise_for_status()
        data = resp.json()
        return data.get("response", "").strip()
    except Exception as e:
        logger.error(f"Ollama chat failed: {e}")
        return ""


PROJECTS_BASE = "projects"


# Default Role Rates (USD/month)
ROLE_RATE_MAP: Dict[str, float] = {
    "Backend Developer": 3000.0,
    "Frontend Developer": 2800.0,
    "QA Analyst": 1800.0,
    "QA Engineer": 2000.0,
    "Data Engineer": 2800.0,
    "Data Analyst": 2200.0,
    "Data Architect": 3500.0,
    "UX Designer": 2500.0,
    "UI/UX Designer": 2600.0,
    "Project Manager": 3500.0,
    "Cloud Engineer": 3000.0,
    "BI Developer": 2700.0,
    "DevOps Engineer": 3200.0,
    "Security Administrator": 3000.0,
    "System Administrator": 2800.0,
    "Solution Architect": 4000.0,
}

#  helpers
def _strip_code_fences(s: str) -> str:
    m = re.search(r"```(?:json)?(.*?)```", s, flags=re.DOTALL | re.IGNORECASE)
    return m.group(1) if m else s

def _extract_json(s: str) -> dict:
    raw = _strip_code_fences(s or "")
    try:
        return json.loads(raw.strip())
    except Exception:
        start, end = raw.find("{"), raw.rfind("}")
        if start >= 0 and end > start:
            try:
                return json.loads(raw[start:end+1])
            except Exception:
                return {}
        return {}
    


def _parse_date_safe(val: Any, fallback: datetime = None) -> datetime:
    """Try to parse a date string; return fallback if invalid."""
    if not val:
        return fallback
    try:
        return datetime.strptime(str(val), "%Y-%m-%d")
    except Exception:
        return fallback

def _safe_str(val: Any) -> str:
    return str(val).strip() if val is not None else ""

async def get_rate_map_for_project(db: AsyncSession, project) -> Dict[str, float]:
    """
    Fetch rate cards for the given project/company.
    Falls back to Sigmoid default rates if none exist
    """
    try:
        # If project has company_id, try fetching company-specific rate cards
        if getattr(project, "company_id", None):
            result = await db.execute(
                select(models.RateCard)
                .filter(models.RateCard.company_id == project.company_id)
            )
            ratecards = result.scalars().all()
            if ratecards:
                return {r.role_name: float(r.monthly_rate) for r in ratecards}

        sigmoid_result = await db.execute(
            select(models.Company).filter(models.Company.name == "Sigmoid")
        )
        sigmoid = sigmoid_result.scalars().first()
        if sigmoid:
            result = await db.execute(
                select(models.RateCard)
                .filter(models.RateCard.company_id == sigmoid.id)
            )
            sigmoid_rates = result.scalars().all()
            if sigmoid_rates:
                return {r.role_name: float(r.monthly_rate) for r in sigmoid_rates}

    except Exception as e:
        logger.warning(f"Failed to fetch rate cards: {e}")
    return ROLE_RATE_MAP


async def _extract_text_from_files(files: List[dict]) -> str:
    results: List[str] = []

    async def _extract_single(f: dict) -> None:
        try:
            blob_bytes = await azure_blob.download_bytes(f["file_path"])
            suffix = os.path.splitext(f["file_name"])[-1].lower()

            def process_file() -> str:
                content = ""
                try:
                    if suffix == ".pdf":
                        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                            tmp.write(blob_bytes)
                            tmp_path = tmp.name
                        try:
                            content = extract_pdf_text(tmp_path)
                        finally:
                            os.remove(tmp_path)

                    elif suffix == ".docx":
                        doc = Document(BytesIO(blob_bytes))
                        content = "\n".join(p.text for p in doc.paragraphs)

                    elif suffix == ".pptx":
                        prs = Presentation(BytesIO(blob_bytes))
                        texts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    texts.append(shape.text)
                        content = "\n".join(texts)

                    elif suffix in [".xlsx", ".xlsm"]:
                        wb = openpyxl.load_workbook(BytesIO(blob_bytes))
                        sheet = wb.active
                        content = "\n".join(
                            " ".join(str(cell) if cell else "" for cell in row)
                            for row in sheet.iter_rows(values_only=True)
                        )

                    elif suffix in [".png", ".jpg", ".jpeg", ".tiff"]:
                        img = Image.open(BytesIO(blob_bytes))
                        content = pytesseract.image_to_string(img)

                    else:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                            tmp.write(blob_bytes)
                            tmp_path = tmp.name
                        try:
                            with open(tmp_path, "r", encoding="utf-8", errors="ignore") as fh:
                                content = fh.read()
                        finally:
                            os.remove(tmp_path)

                except Exception as e:
                    logger.warning(f"Extraction failed for {f['file_name']}: {e}")

                return content.strip()

            text = await anyio.to_thread.run_sync(process_file)

            if text:
                results.append(text)
            else:
                logger.warning(f"Extracted no text from {f['file_name']}")

        except Exception as e:
            logger.warning(f"Failed to extract {f.get('file_name')} (path={f.get('file_path')}): {e}")

    async with anyio.create_task_group() as tg:
        for f in files:
            tg.start_soon(_extract_single, f)

    return "\n\n".join(results)


def _rag_retrieve(query: str, k: int = 5) -> List[Dict]:
    """
    Retrieve semantically similar chunks from Qdrant for RAG.
    Uses Ollama embedding model and returns list of matched chunks.
    Skips retrieval if no valid embedding found.
    """
    try:
        q_emb_list = embed_text_ollama([query])

        # Skip if no valid embeddings returned
        if not q_emb_list or not q_emb_list[0]:
            logger.warning("⚠️ No valid embedding generated — skipping Qdrant retrieval.")
            return []

        q_emb = q_emb_list[0]

        # Sanity check vector dimension
        if not isinstance(q_emb, list) or len(q_emb) == 0:
            logger.warning("⚠️ Empty embedding vector — skipping retrieval.")
            return []

        client = get_qdrant_client()
        results = client.search(
            collection_name=QDRANT_COLLECTION,
            query_vector=q_emb,
            limit=k,
            with_payload=True
        )

        hits = []
        for r in results:
            payload = r.payload or {}
            hits.append({
                "id": payload.get("chunk_id", str(r.id)),
                "parent_id": payload.get("parent_id"),
                "content": payload.get("chunk", ""),
                "title": payload.get("title", ""),
                "score": r.score,
            })

        # Group by parent_id for consistency
        grouped = {}
        for h in hits:
            grouped.setdefault(h["parent_id"], []).append({
                "id": h["id"],
                "content": h["content"],
                "title": h["title"],
                "score": h["score"],
            })

        return [
            {"parent_id": pid, "chunks": chs}
            for pid, chs in grouped.items()
        ]

    except Exception as e:
        logger.warning(f"RAG retrieval (Qdrant) failed: {e}")
        return []

def _build_scope_prompt(rfp_text: str, kb_chunks: List[str], project=None, questions_context: str | None = None) -> str:
    import tiktoken

    # Tokenizer
    tokenizer = tiktoken.get_encoding("cl100k_base")
    # Safe token budget (128k, keep ~4k for completion & system messages)
    context_limit = 128000
    max_total_tokens = context_limit - 4000
    used_tokens = 0

    # Trim RFP text
    rfp_tokens = tokenizer.encode(rfp_text or "")
    if len(rfp_tokens) > 3000:
        rfp_tokens = rfp_tokens[:3000]
    rfp_text = tokenizer.decode(rfp_tokens)
    used_tokens += len(rfp_tokens)

    # Trim KB context
    safe_kb_chunks = []
    for ch in kb_chunks or []:
        tokens = tokenizer.encode(ch)
        if used_tokens + len(tokens) > max_total_tokens:
            break
        safe_kb_chunks.append(ch)
        used_tokens += len(tokens)

    kb_context = "\n\n".join(safe_kb_chunks) if safe_kb_chunks else "(no KB context found)"

    name = (getattr(project, "name", "") or "").strip()
    domain = (getattr(project, "domain", "") or "").strip()
    complexity = (getattr(project, "complexity", "") or "").strip()
    tech_stack = (getattr(project, "tech_stack", "") or "").strip()
    use_cases = (getattr(project, "use_cases", "") or "").strip()
    compliance = (getattr(project, "compliance", "") or "").strip()
    duration = str(getattr(project, "duration", "") or "").strip()

    user_context = (
        "Some overview fields have been provided by the user.\n"
        "Treat these user-provided values as the source of truth.\n"
        "Only fill in fields that are blank — do NOT overwrite the given values.\n\n"
        f"Project Name: {name or '(infer if missing)'}\n"
        f"Domain: {domain or '(infer if missing)'}\n"
        f"Complexity: {complexity or '(infer if missing)'}\n"
        f"Tech Stack: {tech_stack or '(infer if missing)'}\n"
        f"Use Cases: {use_cases or '(infer if missing)'}\n"
        f"Compliance: {compliance or '(infer if missing)'}\n"
        f"Duration (months): {duration or '(infer if missing)'}\n\n"
    )

    today_str = datetime.today().date().isoformat()

    return (
        "You are an expert AI project planner.\n"
        "Use the RFP/project text as the **primary source** \n"
        "Use questions and answers to clarify ambiguities.\n"
        "but enrich missing fields with the Knowledge Base context (if relevant).\n"
        "Return ONLY valid JSON (no prose, no markdown, no commentary).\n\n"
        "Output schema:\n"
        "{\n"
        '  "overview": {\n'
        '    "Project Name": string,\n'
        '    "Domain": string,\n'
        '    "Complexity": string,\n'
        '    "Tech Stack": string,\n'
        '    "Use Cases": string,\n'
        '    "Compliance": string,\n'
        '    "Duration": number\n'
        "  },\n"
        '  "activities": [\n'
        '    {\n'
        '      "ID": int,\n'
        '      "Activities": string,\n'
        '      "Description": string | null,\n'
        '      "Owner": string | null,\n'
        '      "Resources": string | null,\n'
        '      "Start Date": "yyyy-mm-dd",\n'
        '      "End Date": "yyyy-mm-dd",\n'
        '      "Effort Months": number\n'
        "    }\n"
        "  ],\n"
        '  "resourcing_plan": [],\n'
        '  "project_summary": {\n'
        '    "executive_summary": string (2-3 paragraphs overview),\n'
        '    "key_deliverables": [string] (list of 5-7 main deliverables),\n'
        '    "success_criteria": [string] (list of 3-5 success metrics),\n'
        '    "risks_and_mitigation": [{"risk": string, "mitigation": string}] (3-4 key risks)\n'
        "  }\n"
        "}\n\n"
        "Scheduling Rules: \n"
        f"- The first activity must always start today ({today_str}).\n"
        "- If two activities are **independent**, overlap their timelines by **70–80%** of their duration (not full overlap)."
        "- If one activity **depends** on another, allow a small overlap of **10-15%** near the end of the predecessor if feasible."
        "- Avoid full serialization unless strictly required by dependency."
        "- Avoid full parallelism where all tasks start together — stagger independent ones by **5-10%**."
        "- Ensure overall project duration stays **≤ 12 months**."
        "- Auto-calculate **End Date = Start Date + Effort Months**.\n"
        "- Auto-calculate **overview.Duration** as the total span in months from the earliest Start Date to the latest End Date.\n"
        "- `Complexity` should be simple, medium, or large based on duration of project.\n"
        "- **Always assign at least one Resource**."
        "- Distinguish `Owner` (responsible lead role) and `Resources` (supporting roles)."
        "- `Owner` and `Resources` must be valid IT roles (e.g., Backend Developer, AI Engineer, QA Engineer, etc.)."
        "- `Owner` is always a role who manages that particular activity (not a personal name).\n"
        "- `Resources` must contain only roles which are required for that particular activity, distinct from `Owner`.\n"
        "- If `Resources` is missing, fallback to the same `Owner` role.\n"
        "- Use less resources as much as possible.\n"
        "- Effort Months should be small numbers 0.5 to 1.5 months (inclusive).\n"
        "- IDs must start from 1 and increment sequentially.\n"
        "- If the RFP or Knowledge Base text lacks detail, infer the missing pieces logically."
        "- Include all relevant roles and activities that ensure delivery of the project scope."
        "- Keep all field names exactly as in the schema.\n"
        "- Generate a comprehensive project_summary with:\n"
        "  * executive_summary: 2-3 paragraph high-level overview of the project, objectives, and expected outcomes\n"
        "  * key_deliverables: List 5-7 concrete deliverables (e.g., 'Production-ready web application', 'API documentation', etc.)\n"
        "  * success_criteria: List 3-5 measurable success metrics (e.g., '99.9% uptime', 'Response time < 200ms', etc.)\n"
        "  * risks_and_mitigation: List 3-4 key risks with mitigation strategies (e.g., risk: 'Third-party API dependency', mitigation: 'Implement fallback mechanisms')\n"
        f"{user_context}"
        f"RFP / Project Files Content:\n{rfp_text}\n\n"
        f"Knowledge Base Context (for enrichment only):\n{kb_context}\n"
        f"Clarification Q&A (User-confirmed answers take highest priority)\n"
        f"Use these answers to override or clarify any ambiguous or conflicting information.\n"
        f"Do NOT hallucinate beyond these facts.\n\n"
        f"{questions_context}\n"
    )


def _build_questionnaire_prompt(rfp_text: str, kb_chunks: List[str], project=None) -> str:
    """
    Build a prompt that forces the model to infer categories dynamically from RFP context.
    """
    name = getattr(project, "name", "Unnamed Project")
    domain = getattr(project, "domain", "General")
    tech = getattr(project, "tech_stack", "Modern Web Stack")
    compliance = getattr(project, "compliance", "General")
    duration = getattr(project, "duration", "TBD")

    return f"""
You are a **senior business analyst** preparing a requirement-clarification questionnaire
based on an RFP document.

Your goal: identify the main THEMES and subareas discussed in the RFP or Knowledge Base,
and then create **categories of questions** that align with those themes.
Do NOT reuse example categories blindly — derive them from the content itself.

---

### Project Context
- Project Name: {name}
- Domain: {domain}
- Tech Stack: {tech}
- Compliance: {compliance}
- Duration: {duration}

### RFP Content
{rfp_text}

### Knowledge Base Context
{kb_chunks}

---

### TASK
1. First, analyze the RFP text to identify **key themes or topics** (e.g., Data Governance, SOX Controls,
   Cloud Migration, AI Enablement, Supply Chain Optimization, etc.).
2. For each theme, create a **category** with 5-6 specific questions.
3. Questions should clarify requirements, assumptions, or current-state processes.
4. Avoid repeating generic categories like "Architecture" or "Data & Security"
   unless they are explicitly discussed in the RFP.

---

### OUTPUT FORMAT
Return ONLY valid JSON in this structure:

{{
  "questions": [
    {{
      "category": "Data Governance & Ownership",
      "items": [
        {{
          "question": "Is there a defined data ownership model for finance data?",
          "user_understanding": "",
          "comment": ""
        }},
        {{
          "question": "Do you maintain audit logs for data corrections?",
          "user_understanding": "",
          "comment": ""
        }}
      ]
    }},
    {{
      "category": "Regulatory Readiness and SOX Scope",
      "items": [
        {{
          "question": "What parts of the organization are in SOX scope?",
          "user_understanding": "",
          "comment": ""
        }}
      ]
    }}
  ]
}}

### RULES
- Categories must emerge logically from the RFP and KB text.
- Each category must contain at least 2 context-relevant questions.
- Each question must be concise, unambiguous, and require a short descriptive answer.
- Always include empty strings for 'user_understanding' and 'comment'.
- Output ONLY valid JSON (no explanations or markdown).
"""

def _extract_questions_from_text(raw_text: str) -> list[dict]:
    try:
        parsed = _extract_json(raw_text)

        # Case 1: Proper JSON with nested categories
        if isinstance(parsed, dict) and "questions" in parsed:
            qdata = parsed["questions"]
            if isinstance(qdata, list) and all(isinstance(x, dict) for x in qdata):
                # check if already nested structure
                if "items" in qdata[0]:
                    normalized = []
                    for cat in qdata:
                        normalized.append({
                            "category": cat.get("category", "General"),
                            "items": [
                                {
                                    "question": i.get("question", ""),
                                    "user_understanding": i.get("user_understanding", ""),
                                    "comment": i.get("comment", "")
                                } for i in cat.get("items", [])
                            ]
                        })
                    return normalized

                # Otherwise, flat → group by category
                grouped = {}
                for q in qdata:
                    cat = q.get("category", "General") if isinstance(q, dict) else "General"
                    que = q.get("question", q) if isinstance(q, dict) else str(q)
                    grouped.setdefault(cat, []).append({
                        "question": que,
                        "user_understanding": "",
                        "comment": ""
                    })
                return [{"category": c, "items": lst} for c, lst in grouped.items()]

        # Case 2: List of plain questions
        if isinstance(parsed, list):
            return [{
                "category": "General",
                "items": [{"question": str(q), "user_understanding": "", "comment": ""} for q in parsed]
            }]
    except Exception:
        pass

    # Fallback — parse raw text
    current_cat = "General"
    grouped: dict[str, list] = {}
    for line in raw_text.splitlines():
        line = line.strip()
        if not line:
            continue
        if re.match(r"^(#+\s*)?([A-Z][A-Za-z\s&/]+):?$", line) and not line.endswith("?"):
            current_cat = re.sub(r"^#+\s*", "", line).strip(": ").strip()
            continue
        if "?" in line:
            qtext = re.sub(r"^\d+[\).\s]+", "", line).strip()
            grouped.setdefault(current_cat, []).append({
                "question": qtext,
                "user_understanding": "",
                "comment": ""
            })

    return [{"category": c, "items": lst} for c, lst in grouped.items()]
async def generate_project_questions(db: AsyncSession, project) -> dict:
    """
    Generate a categorized questionnaire for the given project using Ollama.
    Saves the questions.json file in Azure Blob.
    """

    # ---------- Extract RFP ----------
    rfp_text = ""
    try:
        if getattr(project, "files", None):
            files = [{"file_name": f.file_name, "file_path": f.file_path} for f in project.files]
            if files:
                rfp_text = await _extract_text_from_files(files)
    except Exception as e:
        logger.warning(f"Failed to extract RFP for questions: {e}")

    # ---------- Retrieve Knowledge Base ----------
    kb_results = _rag_retrieve(rfp_text or project.name or project.domain)
    kb_chunks = [ch["content"] for group in kb_results for ch in group["chunks"]] if kb_results else []

    # ---------- Build prompt ----------
    prompt = _build_questionnaire_prompt(rfp_text, kb_chunks, project)

    # ---------- Query Ollama ----------
    try:
        raw_text = await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, temperature=0.8))
        questions = _extract_questions_from_text(raw_text)
        total_q = sum(len(cat["items"]) for cat in questions)
        logger.info(f" Generated {total_q} questions under {len(questions)} categories for project {project.id}")

        # ---------- Save to Blob Storage ----------
        blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
        try:
            await azure_blob.upload_bytes(
                json.dumps({"questions": questions}, ensure_ascii=False, indent=2).encode("utf-8"),
                blob_name,
            )

            db_file = models.ProjectFile(
                project_id=project.id,
                file_name="questions.json",
                file_path=blob_name,
            )

            db.add(db_file)
            await db.commit()
            await db.refresh(db_file)

            logger.info(f" Saved questions.json for project {project.id}")
        except Exception as e:
            logger.warning(f"Failed to upload questions.json: {e}")

        return {"questions": questions}

    except Exception as e:
        logger.error(f" Question generation failed: {e}")
        return {"questions": []}
    
# Update questions.json with user input answers
async def update_questions_with_user_input(
    db: AsyncSession, project, user_answers: dict
) -> dict:
    from app.utils import azure_blob

    blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
    try:
        # Load current questions.json
        q_bytes = await azure_blob.download_bytes(blob_name)
        q_json = json.loads(q_bytes.decode("utf-8"))
        questions = q_json.get("questions", [])

        # Merge answers into the structure
        for cat in questions:
            cat_name = cat.get("category")
            for item in cat.get("items", []):
                q_text = item.get("question")
                ans = (
                    user_answers.get(cat_name, {}).get(q_text)
                    if user_answers.get(cat_name)
                    else None
                )
                if ans:
                    item["user_understanding"] = ans

        # Upload updated JSON to Blob
        new_bytes = json.dumps({"questions": questions}, ensure_ascii=False, indent=2).encode("utf-8")
        await azure_blob.upload_bytes(new_bytes, blob_name)
        logger.info(f" Updated questions.json with user input for project {project.id}")

        #  Save / update DB record
        db_file = models.ProjectFile(
            project_id=project.id,
            file_name="questions.json",
            file_path=blob_name,
        )
        db.add(db_file)
        await db.commit()
        await db.refresh(db_file)

        return {"questions": questions}

    except Exception as e:
        logger.error(f"Failed to update questions.json with user input: {e}")
        return {}

    
def _build_architecture_prompt(rfp_text: str, kb_chunks: List[str], project=None) -> str:
    name = (getattr(project, "name", "") or "Untitled Project").strip()
    domain = (getattr(project, "domain", "") or "General").strip()
    tech = (getattr(project, "tech_stack", "") or "Modern Web + Cloud Stack").strip()

    return f"""
    You are a **senior enterprise solution architect** tasked with designing a *tailored cloud system architecture diagram*
    strictly based on the provided RFP and contextual knowledge.

    ### PROJECT CONTEXT
    - **Project Name:** {name}
    - **Domain:** {domain}
    - **Tech Stack:** {tech}

    ### RFP SUMMARY
    {rfp_text}

    ### KNOWLEDGE BASE CONTEXT
    {kb_chunks}

    ---

    ###  STEP 1 — Reasoning (Internal)
    Analyze the provided RFP and knowledge base to:
    1. Identify all domain-specific **entities, systems, or technologies** mentioned or implied.
    2. Categorize each component into the most appropriate architecture layer:
    - Frontend (UI/Apps)
    - Backend (Services/APIs)
    - Data (Databases, Storage, External APIs)
    - AI/Analytics (ML, Insights, NLP, Recommendations)
    - Security/Monitoring/DevOps (IAM, Key Vault, CI/CD, Logging)
    3. Infer **connections and data flows** between components (e.g., API requests, pipelines, message queues).
    4. Skip any layers not relevant to this RFP.

    You will use this reasoning to build the architecture — but **do not include this reasoning** in your final output.

    ---

    ###  STEP 2 — Graphviz DOT Output
    Generate **only valid Graphviz DOT code** representing the inferred architecture.

    Follow these rules strictly:
    - Begin with: `digraph Architecture {{`
    - End with: `}}`
    - Use **horizontal layout** → `rankdir=LR`
    - Include **only relevant clusters** (omit unused layers)
    - Keep ≤ 15 nodes total
    - Use **orthogonal edges** (`splines=ortho`)
    - Each node label must clearly represent an actual system, service, or tool
    - Logical flow should follow Frontend → Backend → Data → AI → Security (only if applicable)
    -  **Ensure data layers both receive and provide information** — show arrows *into* and *out of* data/storage nodes if analytics, AI, or reporting components exist.

    ---

    ### VISUAL STYLE
    - **Graph:** dpi=200, bgcolor="white", nodesep=1.3, ranksep=1.3
    - **Clusters:** style="filled,rounded", fontname="Helvetica-Bold", fontsize=13
    - **Node Shapes and Colors:**
    - Frontend → `box`, pastel blue (`fillcolor="#E3F2FD"`)
    - Backend/API → `box3d`, pastel green (`fillcolor="#E8F5E9"`)
    - Data/Storage → `cylinder`, pastel yellow (`fillcolor="#FFFDE7"`)
    - AI/Analytics → `ellipse`, pastel purple (`fillcolor="#F3E5F5"`)
    - Security/Monitoring → `diamond`, gray (`fillcolor="#ECEFF1"`)
    - **Edges:** color="#607D8B", penwidth=1.5, arrowsize=0.9

    ---

    ###  STEP 3 — Domain Intelligence (Auto-Enrichment)
    If applicable, automatically enrich the architecture using these domain patterns:

    - **FinTech** → Payment Gateway, Fraud Detection, KYC/AML Service, Ledger DB
    - **HealthTech** → Patient Portal, EHR System, FHIR API, HIPAA Compliance Layer
    - **GovTech** → Citizen Portal, Secure API Gateway, Compliance & Audit Logging
    - **AI/ML Projects** → Model API, Embedding Store, Training Pipeline, Monitoring Service
    - **Data Platforms** → ETL Pipeline, Data Lake, BI Dashboard
    - **Enterprise SaaS** → Tenant Manager, Auth Service, Billing & Subscription Module

    Include these elements **only if they logically fit** the RFP description.

    ---

    ###  STEP 4 — OUTPUT RULES
    - Output *only* the Graphviz DOT syntax — **no markdown**, **no reasoning**, **no commentary**
    - The final response should be a single valid DOT diagram ready for rendering
    """

async def _generate_fallback_architecture(
    db: AsyncSession,
    project,
    blob_base_path: str
) -> tuple[models.ProjectFile | None, str]:
    """
    Generate and upload a default fallback architecture diagram (4-layer generic layout).
    Triggered when Ollama or Graphviz generation fails.
    """
    logger.warning(" Using fallback default architecture layout")

    # --- Default DOT diagram ---
    fallback_dot = """
digraph Architecture {
    rankdir=LR;
    graph [dpi=200, bgcolor="white", nodesep=1.3, ranksep=1.2, splines=ortho];
    node [style="rounded,filled", fontname="Helvetica-Bold", fontsize=13, penwidth=1.2];

    subgraph cluster_frontend {
        label="Frontend / User Touchpoints";
        style="filled,rounded"; fillcolor="#E3F2FD";
        web[label="Web App (React / Angular)", shape=box, fillcolor="#BBDEFB"];
        mobile[label="Mobile App", shape=box, fillcolor="#BBDEFB"];
    }

    subgraph cluster_backend {
        label="Backend / Services";
        style="filled,rounded"; fillcolor="#E8F5E9";
        api[label="Core API (FastAPI / Node.js)", shape=box3d, fillcolor="#C8E6C9"];
        auth[label="Auth Service", shape=box3d, fillcolor="#C8E6C9"];
    }

    subgraph cluster_data {
        label="Data / Storage";
        style="filled,rounded"; fillcolor="#FFFDE7";
        db[label="Database (PostgreSQL)", shape=cylinder, fillcolor="#FFF9C4"];
        blob[label="Blob Storage", shape=cylinder, fillcolor="#FFF9C4"];
    }

    subgraph cluster_ai {
        label="AI / Analytics";
        style="filled,rounded"; fillcolor="#F3E5F5";
        ai[label="AI Engine / Insights", shape=ellipse, fillcolor="#E1BEE7"];
        dashboard[label="BI Dashboard", shape=ellipse, fillcolor="#E1BEE7"];
    }

    # Data flow (using xlabels to avoid orthogonal label warnings)
    web -> api [xlabel="HTTP Request"];
    mobile -> api [xlabel="Mobile API Call"];
    api -> db [xlabel="DB Query"];
    db -> ai [xlabel="ETL/Inference"];
    ai -> dashboard [xlabel="Visualization"];
    api -> auth [xlabel="Auth Validation"];

}
"""

    # --- Render DOT → PNG & SVG ---
    tmp_base = tempfile.NamedTemporaryFile(delete=False, suffix=".dot").name
    try:
        graph = graphviz.Source(fallback_dot, engine="dot")
        graph.render(tmp_base, format="png", cleanup=True)
        graph.render(tmp_base, format="svg", cleanup=True)

        png_path = tmp_base + ".png"
        svg_path = tmp_base + ".svg"
    except Exception as e:
        logger.error(f" Fallback Graphviz rendering failed: {e}")
        return None, ""

    # --- Upload both files to Azure Blob ---
    blob_name_png = f"{blob_base_path}/architecture_fallback_{project.id}.png"
    blob_name_svg = f"{blob_base_path}/architecture_fallback_{project.id}.svg"

    try:
        with open(png_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_png)
        with open(svg_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_svg)
    finally:
        for path in [png_path, svg_path, tmp_base]:
            try:
                os.remove(path)
            except FileNotFoundError:
                pass

    # --- Save both records in DB ---
    db_file_png = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.png",
        file_path=blob_name_png,
    )
    db_file_svg = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.svg",
        file_path=blob_name_svg,
    )

    db.add_all([db_file_png, db_file_svg])
    await db.commit()
    await db.refresh(db_file_png)
    await db.refresh(db_file_svg)

    logger.info(
        f" Fallback architecture diagrams stored for project {project.id}: "
        f"{blob_name_png}, {blob_name_svg}"
    )

    return db_file_png, blob_name_png



async def generate_architecture(
    db: AsyncSession,
    project,
    rfp_text: str,
    kb_chunks: List[str],
    blob_base_path: str,
) -> tuple[models.ProjectFile | None, str]:
    """
    Generate a visually clean, context-aware architecture diagram (PNG & SVG)
    from RFP + KB context using Ollama + Graphviz.
    Uses dynamic prompts that adapt layers automatically (no static template).
    Includes retry logic, sanitization, validation, and fallback diagram.
    """

    prompt = _build_architecture_prompt(rfp_text, kb_chunks, project)

    # ---------- Step 1: Ask Ollama for Graphviz DOT code ----------
    async def _generate_dot_from_ai(retry: int = 0) -> str:
        """Call Ollama locally to generate DOT diagram."""
        try:
            return await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, temperature=0.7))
        except Exception as e:
            if retry < 2:
                logger.warning(f"Ollama call failed (retry {retry+1}/3): {e}")
                await anyio.sleep(2)
                return await _generate_dot_from_ai(retry + 1)
            logger.error(f"Ollama architecture generation failed after retries: {e}")
            return ""


    dot_code = await _generate_dot_from_ai()
    if not dot_code:
        logger.warning(" No DOT code returned by AI — generating fallback diagram")
        return await _generate_fallback_architecture(db, project, blob_base_path)

    # ---------- Step 2: Clean & sanitize DOT ----------
    dot_code = re.sub(r"```[a-zA-Z]*", "", dot_code).replace("```", "").strip()
    dot_code = dot_code.strip("`").strip()
    dot_code = re.sub(r"(?i)^graph\s", "digraph ", dot_code)

    # Fix brace mismatch
    open_braces = dot_code.count("{")
    close_braces = dot_code.count("}")
    if open_braces > close_braces:
        dot_code += "}" * (open_braces - close_braces)
    elif close_braces > open_braces:
        dot_code = "digraph Architecture {\n" + dot_code

    if not dot_code.lower().startswith("digraph"):
        dot_code = f"digraph Architecture {{\n{dot_code}\n}}"

    # Remove control characters
    dot_code = re.sub(r"[^\x09\x0A\x0D\x20-\x7E]", "", dot_code)

    # ---------- Step 3: Do NOT override GPT’s style ----------
    # Keep GPT’s own clusters, nodes, and colors — just ensure it's syntactically valid
    # (Old static preamble removed intentionally)

    # ---------- Step 4: Render DOT → PNG & SVG ----------
    try:
        tmp_base = tempfile.NamedTemporaryFile(delete=False, suffix=".dot").name
        graph = graphviz.Source(dot_code, engine="dot")

        # Render both PNG and SVG for better clarity
        graph.render(tmp_base, format="png", cleanup=True)
        graph.render(tmp_base, format="svg", cleanup=True)

        png_path = tmp_base + ".png"
        svg_path = tmp_base + ".svg"
    except Exception as e:
        logger.error(f" Graphviz rendering failed: {e}\n--- DOT Snippet ---\n{dot_code[:800]}")
        return await _generate_fallback_architecture(db, project, blob_base_path)

    # ---------- Step 5: Upload PNG to Azure Blob ----------
    blob_name_png = f"{blob_base_path}/architecture_{project.id}.png"
    blob_name_svg = f"{blob_base_path}/architecture_{project.id}.svg"

    try:
        with open(png_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_png)

        with open(svg_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_svg)
    finally:
        for path in [png_path, svg_path, tmp_base]:
            try:
                os.remove(path)
            except FileNotFoundError:
                pass

    # ---------- Step 6: Replace old architecture file ----------
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project.id,
            models.ProjectFile.file_name == "architecture.png",
        )
    )
    old_file = result.scalars().first()
    if old_file:
        try:
            await azure_blob.delete_blob(old_file.file_path)
            await db.delete(old_file)
            await db.commit()
        except Exception as e:
            logger.warning(f" Failed to delete old architecture.png: {e}")

    # ---------- Step 7: Save new ProjectFile records (PNG + SVG) ----------
    db_file_png = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.png",
        file_path=blob_name_png,
    )
    db_file_svg = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.svg",
        file_path=blob_name_svg,
    )

    db.add_all([db_file_png, db_file_svg])
    await db.commit()
    await db.refresh(db_file_png)
    await db.refresh(db_file_svg)

    logger.info(
        f" Architecture diagrams stored successfully for project {project.id}: "
        f"{blob_name_png}, {blob_name_svg}"
    )

    return db_file_png, blob_name_png


# --- Cleaner ---
async def clean_scope(db: AsyncSession, data: Dict[str, Any], project=None) -> Dict[str, Any]:
    if not isinstance(data, dict):
        return {}

    ist = pytz.timezone("Asia/Kolkata")
    # Use timezone-naive datetime to avoid comparison issues with parsed dates
    today = datetime.now(ist).replace(hour=0, minute=0, second=0, microsecond=0, tzinfo=None)

    activities: List[Dict[str, Any]] = []
    start_dates, end_dates = [], []
    role_month_map: Dict[str, Dict[str, float]] = {}
    role_order: List[str] = []

    # --- Helper: compute monthly allocation based on actual days in month ---
    def month_effort(s: datetime, e: datetime) -> Dict[str, float]:
        cur = s
        month_eff = {}
        while cur <= e:
            year, month = cur.year, cur.month
            days_in_month = monthrange(year, month)[1]
            start_day = cur.day if cur.month == s.month else 1
            end_day = e.day if cur.month == e.month else days_in_month
            days_count = end_day - start_day + 1
            month_eff[f"{cur.strftime('%b %Y')}"] = round(days_count / 30.0, 2)
            # move to next month
            if month == 12:
                cur = datetime(year + 1, 1, 1)
            else:
                cur = datetime(cur.year, cur.month + 1, 1)
        return month_eff

    # --- Process activities ---
    for idx, a in enumerate(data.get("activities") or [], start=1):
        owner = a.get("Owner") or "Unassigned"

        # Parse dependencies
        raw_deps = [d.strip() for d in str(a.get("Resources") or "").split(",") if d.strip()]

        # Remove owner from resources if duplicated
        raw_deps = [r for r in raw_deps if r.lower() != owner.lower()]

        # Owner always included, then other resources
        roles = [owner] + raw_deps

        s = _parse_date_safe(a.get("Start Date"), today)
        e = _parse_date_safe(a.get("End Date"), s + timedelta(days=30))
        if e < s:
            e = s + timedelta(days=30)

        # --- allocate per month (no splitting among roles) ---
        month_alloc = month_effort(s, e)
        for role in roles:
            if role not in role_month_map:
                role_month_map[role] = {}
                role_order.append(role)
            for m, eff in month_alloc.items():
                role_month_map[role][m] = role_month_map[role].get(m, 0.0) + eff

        dur_days = max(1, (e - s).days)
        activities.append({
            "ID": idx,
            "Activities": _safe_str(a.get("Activities")),
            "Description": _safe_str(a.get("Description")),
            "Owner": owner,
            "Resources": ", ".join(raw_deps), 
            "Start Date": s,
            "End Date": e,
            "Effort Months": round(dur_days / 30.0, 2),
        })

        start_dates.append(s)
        end_dates.append(e)

        # --- Sort activities ---
    activities.sort(key=lambda x: x["Start Date"])
    for idx, a in enumerate(activities, start=1):
        a["ID"] = idx
        a["Start Date"] = a["Start Date"].strftime("%Y-%m-%d")
        a["End Date"] = a["End Date"].strftime("%Y-%m-%d")

    # --- Project span & month labels (Month 1, Month 2, ...) ---
    min_start = min(start_dates) if start_dates else today
    max_end = max(end_dates) if end_dates else min_start
    duration = max(1.0, round(max(1, (max_end - min_start).days) / 30.0, 2))
    total_months = max(1, math.ceil((max_end - min_start).days / 30.0))

    month_labels = [f"Month {i}" for i in range(1, total_months + 1)]

    # --- Build per-role, per-month day usage ---
    role_month_usage: Dict[str, Dict[str, float]] = {r: {m: 0.0 for m in month_labels} for r in role_order}

    # Compute total active days per relative month window
    for act in activities:
        s = _parse_date_safe(act.get("Start Date"), today)
        e = _parse_date_safe(act.get("End Date"), s + timedelta(days=30))
        if e < s:
            e = s + timedelta(days=30)

        involved_roles = [act.get("Owner") or "Unassigned"] + [
            r.strip() for r in str(act.get("Resources") or "").split(",") if r.strip()
        ]

        for m_idx in range(total_months):
            rel_start = min_start + timedelta(days=m_idx * 30)
            rel_end = min_start + timedelta(days=(m_idx + 1) * 30)

            # overlap between activity and this relative month window
            overlap_start = max(s, rel_start)
            overlap_end = min(e, rel_end)
            overlap_days = 0
            if overlap_end >= overlap_start:
                overlap_days = (overlap_end - overlap_start).days + 1

            if overlap_days > 0:
                for r in involved_roles:
                    if r not in role_month_usage:
                        role_month_usage[r] = {ml: 0.0 for ml in month_labels}
                    role_month_usage[r][f"Month {m_idx + 1}"] += overlap_days

    # --- Convert days to effort with 4-tier partial-month logic ---
    for r, months in role_month_usage.items():
        for m, days in months.items():
            if days > 21:
                months[m] = 1.0
            elif 15 <= days <= 21:
                months[m] = 0.75
            elif 8 <= days < 15:
                months[m] = 0.5
            elif 1 <= days < 8:
                months[m] = 0.25
            else:
                months[m] = 0.0

    try:
        if db:
            ROLE_RATE_MAP_DYNAMIC = await get_rate_map_for_project(db, project)
        else:
            ROLE_RATE_MAP_DYNAMIC = ROLE_RATE_MAP
    except Exception as e:
        logger.warning(f"Rate map fallback due to error: {e}")
        ROLE_RATE_MAP_DYNAMIC = ROLE_RATE_MAP


    # --- Build final resourcing plan ---
    resourcing_plan = []
    for idx, role in enumerate(role_order, start=1):
        month_efforts = role_month_usage.get(role, {m: 0 for m in month_labels})
        total_effort = sum(month_efforts.values())
        rate = ROLE_RATE_MAP_DYNAMIC.get(role, ROLE_RATE_MAP.get(role, 2000.0))
        cost = round(total_effort * rate, 2)
        plan_entry = {
            "ID": idx,
            "Resources": role,
            "Rate/month": rate,
            **month_efforts,
            "Efforts": total_effort,
            "Cost": cost,
        }
        resourcing_plan.append(plan_entry)

    # --- Apply discount if present ---
    discount_percentage = data.get("discount_percentage", 0)
    if discount_percentage and isinstance(discount_percentage, (int, float)) and discount_percentage > 0:
        discount_multiplier = 1 - (discount_percentage / 100.0)
        logger.info(f"💰 Applying {discount_percentage}% discount (multiplier: {discount_multiplier})")

        # Apply discount to all costs in resourcing_plan
        for plan_entry in resourcing_plan:
            original_cost = plan_entry.get("Cost", 0)
            discounted_cost = round(original_cost * discount_multiplier, 2)
            plan_entry["Cost"] = discounted_cost
            logger.info(f"  → {plan_entry['Resources']}: ${original_cost} → ${discounted_cost}")

    # --- Overview ---
    ov = data.get("overview") or {}
    data["overview"] = {
        "Project Name": _safe_str(ov.get("Project Name") or getattr(project, "name", "Untitled Project")),
        "Domain": _safe_str(ov.get("Domain") or getattr(project, "domain", "")),
        "Complexity": _safe_str(ov.get("Complexity") or getattr(project, "complexity", "")),
        "Tech Stack": _safe_str(ov.get("Tech Stack") or getattr(project, "tech_stack", "")),
        "Use Cases": _safe_str(ov.get("Use Cases") or getattr(project, "use_cases", "")),
        "Compliance": _safe_str(ov.get("Compliance") or getattr(project, "compliance", "")),
        "Duration": duration,
        "Generated At": datetime.now(ist).strftime("%Y-%m-%d %H:%M %Z"),
    }
    try:
        if getattr(project, "company", None):
            data["overview"]["Currency"] = getattr(project.company, "currency", "USD")
        else:
            data["overview"]["Currency"] = "USD"
    except Exception:
        data["overview"]["Currency"] = "USD"

    # Add discount to overview if present
    if discount_percentage and isinstance(discount_percentage, (int, float)) and discount_percentage > 0:
        data["overview"]["Discount"] = f"{discount_percentage}%"
        total_cost = sum(plan_entry.get("Cost", 0) for plan_entry in resourcing_plan)
        data["overview"]["Total Cost (After Discount)"] = f"${total_cost:,.2f}"

    data["activities"] = activities
    data["resourcing_plan"] = resourcing_plan

    # Keep discount_percentage in output for reference
    if discount_percentage and isinstance(discount_percentage, (int, float)) and discount_percentage > 0:
        data["discount_percentage"] = discount_percentage

    return data


async def generate_project_scope(db: AsyncSession, project) -> dict:
    """
    Generate project scope + architecture diagram + store architecture in DB + return combined JSON.
    """

    #  Ensure the project has a valid company reference (fallback to Sigmoid)
    if not getattr(project, "company_id", None):
        from app.utils import ratecards
        sigmoid = await ratecards.get_or_create_sigmoid_company(db)
        project.company_id = sigmoid.id
        await db.commit()
        await db.refresh(project)
        logger.info(f"Linked project {project.id} to Sigmoid company as fallback")

    tokenizer = tiktoken.get_encoding("cl100k_base")
    context_limit = 128000
    max_total_tokens = context_limit - 4000
    used_tokens = 0

    # ---------- Extract RFP ----------
    rfp_text = ""
    try:
        files: List[dict] = []
        if getattr(project, "files", None):
            try:
                files = [{"file_name": f.file_name, "file_path": f.file_path} for f in project.files]
            except Exception as e:
                logger.warning(f" Could not access project.files: {e}")
                files = []
        if files:
            rfp_text = await _extract_text_from_files(files)
    except Exception as e:
        logger.warning(f"File extraction for project {getattr(project, 'id', None)} failed: {e}")

    # ---------- Trim RFP text ----------
    rfp_tokens = tokenizer.encode(rfp_text or "")
    if len(rfp_tokens) > 5000:
        rfp_tokens = rfp_tokens[:5000]
    rfp_text = tokenizer.decode(rfp_tokens)
    used_tokens += len(rfp_tokens)

    # ---------- Retrieve KB context ----------
    fallback_fields = [
        getattr(project, "name", None),
        getattr(project, "domain", None),
        getattr(project, "complexity", None),
        getattr(project, "tech_stack", None),
        getattr(project, "use_cases", None),
        getattr(project, "compliance", None),
        str(getattr(project, "duration", "")) if getattr(project, "duration", None) else None,
    ]
    fallback_text = " ".join(f for f in fallback_fields if f and str(f).strip())

    # If completely empty, create a detailed specific prompt instead of returning empty scope
    if not (rfp_text.strip() or fallback_text.strip()):
        logger.warning(f"⚠️ No RFP text or project metadata for project {project.id}. Using detailed generic prompt.")
        fallback_text = """
Project Requirements:
- Project Type: Software Development Project
- Domain: Web Application Development
- Complexity: Medium
- Tech Stack: React, Node.js, PostgreSQL, AWS
- Duration: 6 months
- Team Size: 5-7 people

Project Scope:
Create a comprehensive project plan with the following phases:

1. Requirements & Planning Phase (1 month)
   - Gather and document requirements
   - Create technical specifications
   - Set up project infrastructure
   Owner: Project Manager
   Resources: Business Analyst, Technical Lead

2. Design Phase (1 month)
   - Design system architecture
   - Create UI/UX mockups
   - Database schema design
   Owner: Solution Architect
   Resources: UI/UX Designer, Database Administrator

3. Development Phase (2.5 months)
   - Frontend development (React)
   - Backend API development (Node.js)
   - Database implementation
   - Integration testing
   Owner: Technical Lead
   Resources: Frontend Developer, Backend Developer, QA Engineer

4. Testing & QA Phase (1 month)
   - Unit testing
   - Integration testing
   - User acceptance testing
   - Bug fixes
   Owner: QA Lead
   Resources: QA Engineer, Backend Developer

5. Deployment & Go-Live (0.5 months)
   - Production deployment
   - Performance optimization
   - Documentation
   - Training
   Owner: DevOps Engineer
   Resources: Technical Lead, Backend Developer

Project Summary:
- Executive Summary: This project aims to develop a comprehensive web application using React and Node.js to streamline business operations. The solution will provide an intuitive user interface for data management, real-time analytics, and seamless integration with existing systems. Expected outcomes include improved operational efficiency, reduced manual errors, and enhanced user experience.

- Key Deliverables: Production-ready web application, REST API with comprehensive documentation, PostgreSQL database with optimized schema, AWS cloud infrastructure setup, User documentation and training materials, Automated testing suite, Performance monitoring dashboard

- Success Criteria: 99.5% application uptime, Page load time under 2 seconds, Support for 1000+ concurrent users, Zero critical security vulnerabilities, 95% user satisfaction score

- Risks: Third-party API downtime (Mitigation: Implement caching and fallback mechanisms), Database performance bottlenecks (Mitigation: Implement proper indexing and query optimization), Resource availability constraints (Mitigation: Cross-train team members and maintain documentation)

Generate activities with realistic start/end dates, proper role assignments, meaningful descriptions, and a comprehensive project summary.
"""

    kb_results = _rag_retrieve(rfp_text or fallback_text)
    kb_chunks = []
    stop = False
    for group in kb_results:
        for ch in group["chunks"]:
            chunk_tokens = len(tokenizer.encode(ch["content"]))
            if used_tokens + chunk_tokens > max_total_tokens:
                stop = True
                break
            kb_chunks.append(ch["content"])
            used_tokens += chunk_tokens
        if stop:
            break

    logger.info(
        f"Final RFP tokens: {len(rfp_tokens)}, KB tokens: {used_tokens - len(rfp_tokens)}, Total: {used_tokens}/{max_total_tokens}"
    )

    # ---------- Load questions.json (if exists) and build Q&A context ----------
    questions_context = None
    try:
        q_blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
        if await azure_blob.blob_exists(q_blob_name):
            q_bytes = await azure_blob.download_bytes(q_blob_name)
            q_json = json.loads(q_bytes.decode("utf-8"))

            q_lines = []
            for category in q_json.get("questions", []):
                cat_name = category.get("category", "General")
                q_lines.append(f"### {cat_name}")
                for item in category.get("items", []):
                    q = item.get("question", "").strip()
                    a = item.get("user_understanding", "").strip() or "(unanswered)"
                    comment = item.get("comment", "").strip()
                    line = f"Q: {q}\nA: {a}"
                    if comment:
                        line += f"\nComment: {comment}"
                    q_lines.append(line)

            questions_context = "\n".join(q_lines)
            logger.info(f"Loaded {len(q_lines)} question lines for project {project.id}")
        else:
            logger.info(f"No questions.json found for project {project.id}, skipping Q&A context.")

    except Exception as e:
        logger.warning(f" Could not include questions.json context: {e}")
        questions_context = None

    


    # ---------- Build + query ----------
    prompt = _build_scope_prompt(rfp_text, kb_chunks, project, questions_context=questions_context)
    try:
        # Step 1: Generate scope via Ollama
        logger.info(f"🤖 Calling Ollama for scope generation... (prompt length: {len(prompt)} chars)")
        raw_text = await anyio.to_thread.run_sync(lambda: ollama_chat(prompt))
        logger.info(f"📝 Ollama raw response length: {len(raw_text)} chars")
        logger.debug(f"📝 Ollama response preview (first 500 chars): {raw_text[:500]}")

        if not raw_text or len(raw_text.strip()) < 50:
            logger.error(f"❌ Ollama returned empty or too short response: {len(raw_text)} chars")
            logger.error("   This usually means:")
            logger.error("   1. Ollama service is not running properly")
            logger.error("   2. The model (deepseek-r1) is not loaded")
            logger.error("   3. Out of memory or timeout")
            return {}

        raw = _extract_json(raw_text)

        # Validate that LLM actually generated content, not just structure
        if raw.get('activities'):
            activities = raw.get('activities', [])
            empty_fields_count = 0
            for act in activities:
                if (not act.get('Activities', '').strip() or
                    not act.get('Description', '').strip() or
                    act.get('Owner', '').lower() in ['unassigned', '']):
                    empty_fields_count += 1

            if empty_fields_count > len(activities) * 0.7:  # More than 70% are garbage
                logger.error(f"❌ LLM returned {empty_fields_count}/{len(activities)} activities with empty/invalid content!")
                logger.error("   This means Ollama generated JSON structure but NO actual content.")
                logger.error("   Check if:")
                logger.error("   1. Ollama service is running: curl http://localhost:11434/api/tags")
                logger.error("   2. Model is loaded: ollama list")
                logger.error("   3. Sufficient memory available")
                return {}

        cleaned_scope = await clean_scope(db, raw, project=project)
        # Update project fields from generated overview (just like finalize_scope)
        overview = cleaned_scope.get("overview", {})
        if overview:
            project.name = overview.get("Project Name") or project.name
            project.domain = overview.get("Domain") or project.domain
            project.complexity = overview.get("Complexity") or project.complexity
            project.tech_stack = overview.get("Tech Stack") or project.tech_stack
            project.use_cases = overview.get("Use Cases") or project.use_cases
            project.compliance = overview.get("Compliance") or project.compliance
            project.duration = str(overview.get("Duration") or project.duration)

            try:
                await db.commit()
                await db.refresh(project)
                logger.info(f" Project metadata updated from generated scope for project {project.id}")
            except Exception as e:
                logger.warning(f" Failed to update project metadata: {e}")


        # Step 2: Generate + store architecture diagram
        try:
            blob_base_path = f"{PROJECTS_BASE}/{getattr(project, 'id', 'unknown')}"
            db_file, arch_blob = await generate_architecture(
                db, project, rfp_text, kb_chunks, blob_base_path
            )
            cleaned_scope["architecture_diagram"] = arch_blob or None
        except Exception as e:
            logger.warning(f"Architecture diagram generation failed: {e}")
            cleaned_scope["architecture_diagram"] = None

        # Step 3: Auto-save finalized_scope.json in Azure Blob + DB
        try:
            from sqlalchemy import select
            result = await db.execute(
                select(models.ProjectFile).filter(
                    models.ProjectFile.project_id == project.id,
                    models.ProjectFile.file_name == "finalized_scope.json",
                )
            )
            old_file = result.scalars().first()
            if old_file:
                logger.info(f"Overwriting existing finalized_scope.json for project {project.id}")
            else:
                old_file = models.ProjectFile(
                    project_id=project.id,
                    file_name="finalized_scope.json",
                )

            blob_name = f"{PROJECTS_BASE}/{project.id}/finalized_scope.json"

            await azure_blob.upload_bytes(
                json.dumps(cleaned_scope, ensure_ascii=False, indent=2).encode("utf-8"),
                blob_name,
                overwrite=True, 
            )

            old_file.file_path = blob_name
            db.add(old_file)
            await db.commit()
            await db.refresh(old_file)

            logger.info(f" finalized_scope.json overwritten for project {project.id}")

        except Exception as e:
            logger.warning(f" Failed to auto-save finalized_scope.json: {e}")
        return cleaned_scope

    except Exception as e:
        logger.error(f"Ollama scope generation failed: {e}")
        return {}


async def regenerate_from_instructions(
    db: AsyncSession,
    project: models.Project,
    draft: dict,
    instructions: str
) -> dict:
    """
    Regenerate the project scope from user instructions using a creative AI-guided prompt.
    Enhances activity sequencing, roles, and effort estimates while preserving valid JSON structure.
    """
    logger.info(f" Regenerating scope for project {project.id} with creative AI response...")

    if not instructions or not instructions.strip():
        cleaned = await clean_scope(db, draft, project=project)
        return {**cleaned, "_finalized": True}


    prompt = f"""
You are an **expert AI project planner and delivery architect** responsible for maintaining a project scope in JSON format.

You are given:
1. The current draft project scope (JSON with keys: `overview`, `activities`, `resourcing_plan`).
2. The user’s latest change instructions.

Your task:
- **Understand** the user’s intent (instructions may be in natural language).
- **Regenerate** the scope accordingly:
  - Apply all user instructions faithfully.
  - Preserve structure and realism of the plan.
  - Re-calculate activity dates, dependencies, and efforts using the rules below.
  - Reflect improvements like “optimize”, “simplify”, “rebalance”, or “add QA phase”.

---

### RULES OF MODIFICATION

####  Schema
- Preserve the same top-level keys: `overview`, `activities`, `resourcing_plan`.
- Every activity must have: "ID", "Activities", "Description", "Owner", "Resources",
- "Start Date", "End Date", "Effort Months"
- Use valid ISO dates (`yyyy-mm-dd`).
- Keep total duration ≤ 12 months.

**CRITICAL: What activities look like**
CORRECT activity example:
```json
{{
  "ID": 1,
  "Activities": "Project Initiation and Requirements Gathering",
  "Description": "Define project scope, gather requirements, create initial documentation",
  "Owner": "Project Manager",
  "Resources": "Business Analyst, Data Architect",
  "Start Date": "2025-01-15",
  "End Date": "2025-02-28",
  "Effort Months": 1.5
}}
```

WRONG activity example (DO NOT DO THIS):
```json
{{
  "ID": 1,
  "Activities": "Project Manager",  ← WRONG! This is a role name, not an activity!
  "Description": "",  ← WRONG! Must have meaningful description!
  "Owner": "Unassigned",  ← WRONG! Must have a real owner!
  "Resources": "",
  "Start Date": "2025-01-15",
  "End Date": "2025-02-15",
  "Effort Months": 1
}}
```

####  Temporal Adjustment Rules
Use these to keep the schedule consistent and continuous.

**Add new activity (bottom)**  
- Append at the end.  
- Start date = 10 days *before* the current latest end_date.  
- End date = start_date + duration derived from effort_days.  
- Allow small overlap (10-15 %) with the last activity to maximize parallelism.

**Add new activity (in middle)**  
- Insert between the target activities without disturbing global schedule.  
- Preceding activity’s end date remains fixed.  
- Following activity’s start shifts minimally to maintain continuity.  
- Only local dates adjust; efforts remain unchanged for following activities.

**Delete activity**  
- Remove it completely.  
- Do not introduce gaps; subsequent activities retain start/end dates.

**Split activity into two**  
- Divide one activity into two consecutive ones.  
- Combined effort_days = original.  
- Combined duration = original.  
- Other activities’ dates stay the same.

**Merge two activities**
- Combine both into one.
- start_date = min(start of both)
- end_date = max(end of both)
- effort_days = sum(efforts of both)

####  Role Management Rules
Critical: When user requests to add or remove roles, you MUST update BOTH activities and resourcing_plan.

**IMPORTANT: All changes are INCREMENTAL - preserve existing activities unless explicitly deleted!**

**Remove a role (e.g., "remove Business Analyst")**:
1. Keep ALL existing activities
2. Find all activities where the role is the Owner
3. Reassign those activities to another appropriate role
4. Remove the role from ALL Resources fields across all activities
5. DO NOT delete any activities - only change role assignments
6. Example: If removing "Business Analyst":
   - Activity: "Owner": "Business Analyst" → change to "Owner": "Product Manager"
   - Activity: "Resources": "Business Analyst, Data Engineer" → change to "Resources": "Data Engineer"
   - Keep ALL other activities unchanged
   - resourcing_plan: will be auto-calculated

**Add more of an existing role (e.g., "add 1 more Backend Developer")**:
1. **CRITICAL**: Keep ALL existing activities and roles
2. "Add 1 more" means INCREASE allocation, not replace
3. To increase Backend Developer allocation:
   - Add "Backend Developer" to Resources field of MORE existing activities
   - OR extend date ranges of activities that already have Backend Developer
   - OR create 1-2 NEW activities specifically for Backend Developer
4. **DO NOT remove any existing activities or roles**
5. Example: If you have 10 activities and "add 1 Backend Developer":
   - Original: 10 activities with Backend Developer in 3 of them
   - After: Same 10 activities PLUS Backend Developer added to 2-3 more activities
   - Result: Backend Developer effort increases from 3 months to 5-6 months

**Add a new role type (e.g., "add Security Engineer")**:
1. **CRITICAL**: Keep ALL existing activities and roles
2. Add new activities for this role OR add to Resources field of existing activities
3. DO NOT remove any existing activities
4. The resourcing_plan will be auto-generated based on activities

####  Discount Rules
When user requests a discount (e.g., "apply 5% discount", "give 10% discount"):
1. **DO NOT change activities, dates, or efforts**
2. **ONLY note the discount percentage in a special field**
3. Add a new field: "discount_percentage": <number> (e.g., 5 for 5%, 10 for 10%)
4. The discount will be applied automatically during cost calculation
5. Keep all activities, roles, and resourcing_plan calculations unchanged

### Scheduling Rules
- Activities should follow **semi-parallel execution** — overlap realistically but maintain logical order.
- If two activities are **independent**, overlap their timelines by **70–80%** of their duration (not full overlap).
- If one activity **depends** on another, allow a small overlap of **10-15%** near the end of the predecessor if feasible.
- Avoid full serialization unless strictly required by dependency.
- Avoid full parallelism where all tasks start together — stagger independent ones by **10-15%**.
- Ensure overall project duration stays **≤ 12 months**.
- The first activity must always start today.
---

### Regeneration Logic
- Clean and re-order activities logically.
- Maintain coherent dependencies and sequential flow.
- Adjust `overview.duration_months` automatically based on new total project span.
- Keep resource roles realistic and consistent with activities (Backend Developer, Data Engineer, QA Analyst, etc.).
- Reflect optimization or simplification requests (e.g., reduce redundant steps, consolidate phases).

---

###  Output Rules
- Output **only valid JSON** — no markdown, no explanations, no reasoning.
- Must include:
  - `overview` → Project metadata (name, domain, complexity, tech stack, etc.)
  - `activities` → COMPLETE updated list with ALL modifications applied
  - `resourcing_plan` → OPTIONAL (will be auto-calculated from activities)
  - `discount_percentage` → OPTIONAL (only if user requested discount, e.g., 5 for 5%, 10 for 10%)
- **CRITICAL**: If user says "remove [role]", that role MUST NOT appear in ANY activity's Owner or Resources field
- **CRITICAL**: If user says "add 1 more [role]", ADD to existing activities, DO NOT replace them
- **CRITICAL**: If user says "apply X% discount", include "discount_percentage": X in output
- **Dont change schema or field names.**
- **PRESERVE all activities** - only modify/add/remove specific items mentioned by user

---

User Instructions:
{instructions}

Current Draft Scope:
{json.dumps(draft, indent=2, ensure_ascii=False)}

Return only the updated JSON.
"""


    # ---- Query Ollama creatively ----
    # Use lower temperature for more consistent instruction-following
    try:
        raw_text = await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, temperature=0.2))
        logger.info(f"🤖 LLM response length: {len(raw_text)} chars")
        logger.debug(f"LLM raw response (first 500 chars): {raw_text[:500]}")
        updated_scope = _extract_json(raw_text)

        logger.info(f"📊 Extracted scope structure: overview={bool(updated_scope.get('overview'))}, "
                   f"activities={len(updated_scope.get('activities', []))}, "
                   f"resourcing_plan={len(updated_scope.get('resourcing_plan', []))}")

        # Validate activity count - prevent accidental scope replacement
        original_activity_count = len(draft.get('activities', []))
        new_activity_count = len(updated_scope.get('activities', []))
        is_removal_instruction = any(word in instructions.lower() for word in ['remove', 'delete'])

        # Advanced validation: Check if activities are valid/meaningful
        activities_are_valid = True
        validation_failures = []

        if updated_scope.get('activities'):
            unassigned_count = sum(1 for act in updated_scope['activities'] if act.get('Owner', '').lower() in ['unassigned', ''])
            empty_desc_count = sum(1 for act in updated_scope['activities'] if not act.get('Description', '').strip())

            # Check if activity names are just role names (common LLM mistake)
            common_roles = ['project manager', 'business analyst', 'data architect', 'data engineer',
                           'backend developer', 'frontend developer', 'qa engineer', 'devops engineer',
                           'cloud architect', 'data analyst', 'ux designer']
            role_name_activities = sum(1 for act in updated_scope['activities']
                                      if act.get('Activities', '').lower().strip() in common_roles)

            # Check if all activities have identical dates (suspicious)
            dates = [(act.get('Start Date'), act.get('End Date')) for act in updated_scope['activities']]
            unique_date_pairs = len(set(dates))

            # Validation thresholds
            if unassigned_count > new_activity_count * 0.5:  # More than 50% unassigned
                activities_are_valid = False
                validation_failures.append(f"{unassigned_count}/{new_activity_count} activities have Unassigned owner")

            if empty_desc_count > new_activity_count * 0.5:  # More than 50% empty descriptions
                activities_are_valid = False
                validation_failures.append(f"{empty_desc_count}/{new_activity_count} activities have empty descriptions")

            if role_name_activities > new_activity_count * 0.3:  # More than 30% are just role names
                activities_are_valid = False
                validation_failures.append(f"{role_name_activities}/{new_activity_count} activities are named after roles (e.g. 'Project Manager', 'Data Engineer')")

            if unique_date_pairs == 1 and new_activity_count > 1:  # All activities have same dates
                activities_are_valid = False
                validation_failures.append(f"All {new_activity_count} activities have identical dates: {dates[0]}")

        # If LLM significantly reduced activities OR created invalid activities, restore original
        if (new_activity_count < (original_activity_count * 0.7) and not is_removal_instruction) or not activities_are_valid:
            if not activities_are_valid:
                logger.error(f"❌ LLM GENERATED INVALID ACTIVITIES!")
                for failure in validation_failures:
                    logger.error(f"   - {failure}")
            else:
                logger.error(f"❌ LLM LOST TOO MANY ACTIVITIES! Original: {original_activity_count}, New: {new_activity_count}")

            logger.error(f"   User instruction: '{instructions[:100]}'")
            logger.error(f"   🔧 Auto-restoring original activities to prevent data loss")

            # Restore original activities
            updated_scope["activities"] = draft.get("activities", [])
            if "resourcing_plan" not in updated_scope or not updated_scope.get("resourcing_plan"):
                updated_scope["resourcing_plan"] = draft.get("resourcing_plan", [])

            logger.info(f"   ✅ Restored {len(updated_scope['activities'])} valid activities from draft")

        # Log roles found in activities
        if updated_scope.get('activities'):
            owners = set(act.get('Owner', 'Unknown') for act in updated_scope['activities'])
            all_resources = set()
            for act in updated_scope['activities']:
                resources = act.get('Resources', '')
                if resources:
                    all_resources.update(r.strip() for r in str(resources).split(',') if r.strip())
            all_roles = owners | all_resources
            logger.info(f"🎭 Roles in LLM response - Owners: {owners}, Resources: {all_resources}")

            # Validate that "remove" instructions were followed
            if instructions and 'remove' in instructions.lower():
                for role in all_roles:
                    if role.lower() in instructions.lower() and 'remove' in instructions.lower():
                        logger.error(f"❌ LLM FAILED to remove '{role}' - still present in activities despite user instruction!")

            # Validate that "add" instructions were followed
            if instructions and 'add' in instructions.lower():
                # This is harder to validate automatically, but we log for manual inspection
                logger.info(f"ℹ️ User requested to add role(s). Current roles: {all_roles}")

        # Post-processing fallback: manually remove roles if LLM failed
        if instructions and 'remove' in instructions.lower() and updated_scope.get('activities'):
            # Extract role to remove from instructions (basic pattern matching)
            import re
            # Pattern to match "remove <role>" where role can be multi-word
            # Matches everything after "remove" until end of string or common delimiters
            remove_pattern = r'remove\s+([a-zA-Z\s]+?)(?:\s*(?:from|,|\.|\band\b|$))'
            match = re.search(remove_pattern, instructions.lower(), re.IGNORECASE)
            if match:
                role_to_remove = match.group(1).strip()
                logger.info(f"🔧 Post-processing: attempting to remove '{role_to_remove}'")

                # Track if we made changes
                changes_made = False

                # Process each activity
                for act in updated_scope['activities']:
                    # Check if this role is the owner
                    if act.get('Owner', '').lower() == role_to_remove or role_to_remove in act.get('Owner', '').lower():
                        # Find a replacement owner from resources or use a default
                        resources = act.get('Resources', '')
                        if resources and resources.strip():
                            # Use the first resource as the new owner
                            new_owner = resources.split(',')[0].strip()
                            # Remove new owner from resources to avoid duplication
                            remaining_resources = [r.strip() for r in resources.split(',')[1:] if r.strip()]
                            act['Owner'] = new_owner
                            act['Resources'] = ', '.join(remaining_resources)
                            logger.info(f"  → Reassigned activity '{act.get('Activities', 'Unknown')}' from removed role to '{new_owner}'")
                            changes_made = True
                        else:
                            # No resources available, use a generic default
                            act['Owner'] = 'Project Manager'
                            logger.info(f"  → Reassigned activity '{act.get('Activities', 'Unknown')}' from removed role to 'Project Manager'")
                            changes_made = True

                    # Remove from resources field
                    if act.get('Resources'):
                        resources_list = [r.strip() for r in str(act['Resources']).split(',') if r.strip()]
                        # Filter out the role to remove (case-insensitive partial match)
                        filtered_resources = [r for r in resources_list
                                             if role_to_remove not in r.lower() and r.lower() != role_to_remove]
                        if len(filtered_resources) != len(resources_list):
                            act['Resources'] = ', '.join(filtered_resources)
                            changes_made = True

                if changes_made:
                    logger.info(f"✅ Post-processing successfully removed role '{role_to_remove}' from activities")

        # Post-processing: parse discount percentage from instructions
        if instructions:
            import re
            # Pattern to match discount requests: "5% discount", "apply 10% discount", "give 15% discount", etc.
            discount_patterns = [
                r'(\d+)\s*%\s*discount',
                r'discount\s+(?:of\s+)?(\d+)\s*%',
                r'apply\s+(\d+)\s*%',
                r'give\s+(\d+)\s*%',
            ]
            discount_found = False
            for pattern in discount_patterns:
                match = re.search(pattern, instructions.lower())
                if match:
                    discount_percentage = int(match.group(1))
                    logger.info(f"💰 Post-processing: detected {discount_percentage}% discount request")

                    # Add discount to updated_scope if not already present
                    if "discount_percentage" not in updated_scope or not updated_scope.get("discount_percentage"):
                        updated_scope["discount_percentage"] = discount_percentage
                        logger.info(f"  → Added discount_percentage: {discount_percentage}")
                    discount_found = True
                    break

            if not discount_found and any(word in instructions.lower() for word in ['discount', 'reduction', 'reduce cost']):
                logger.warning(f"⚠️ User mentioned discount but couldn't parse percentage. Instructions: {instructions[:100]}")

        # Safety check: if LLM returned empty activities, preserve original
        if not updated_scope.get("activities") or len(updated_scope.get("activities", [])) == 0:
            logger.warning(f"⚠️ LLM returned empty activities - preserving original draft activities")
            logger.info(f"📋 Original draft had {len(draft.get('activities', []))} activities")
            # Preserve original activities and resourcing_plan, but update overview if changed
            updated_scope["activities"] = draft.get("activities", [])
            if "resourcing_plan" not in updated_scope or not updated_scope.get("resourcing_plan"):
                updated_scope["resourcing_plan"] = draft.get("resourcing_plan", [])

        cleaned = await clean_scope(db, updated_scope, project=project)
        logger.info(f"✅ Cleaned scope: {len(cleaned.get('activities', []))} activities, "
                   f"{len(cleaned.get('resourcing_plan', []))} resources")

    except Exception as e:
        logger.error(f" Creative regeneration failed: {e}")
        cleaned = await clean_scope(db, draft, project=project)

    # ---- Update project metadata from overview ----
    overview = cleaned.get("overview", {})
    if overview:
        project.name = overview.get("Project Name") or project.name
        project.domain = overview.get("Domain") or project.domain
        project.complexity = overview.get("Complexity") or project.complexity
        project.tech_stack = overview.get("Tech Stack") or project.tech_stack
        project.use_cases = overview.get("Use Cases") or project.use_cases
        project.compliance = overview.get("Compliance") or project.compliance
        project.duration = str(overview.get("Duration") or project.duration)
        await db.commit()
        await db.refresh(project)
        logger.info(f" Project metadata synced for project {project.id}")

    # ---- Overwrite finalized_scope.json in Blob ----
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project.id,
            models.ProjectFile.file_name == "finalized_scope.json",
        )
    )
    old_file = result.scalars().first() or models.ProjectFile(
        project_id=project.id, file_name="finalized_scope.json"
    )

    blob_name = f"{PROJECTS_BASE}/{project.id}/finalized_scope.json"
    await azure_blob.upload_bytes(
        json.dumps(cleaned, ensure_ascii=False, indent=2).encode("utf-8"),
        blob_name,
        overwrite=True,
    )
    old_file.file_path = blob_name
    db.add(old_file)
    await db.commit()
    await db.refresh(old_file)

    logger.info(f" Creative finalized_scope.json regenerated for project {project.id}")
    return {**cleaned, "_finalized": True}


async def finalize_scope(
    db: AsyncSession,
    project_id: str,
    scope_data: dict
) -> tuple[models.ProjectFile, dict]:
    """
    Finalize the project scope without LLM — just clean, validate sequencing,
    update metadata, and save finalized_scope.json.
    """

    logger.info(f"Finalizing scope (no LLM) for project {project_id}...")

    # ---- Load project ----
    result = await db.execute(
        select(models.Project)
        .options(selectinload(models.Project.company))
        .filter(models.Project.id == project_id)
    )
    project = result.scalars().first()
    if not project:
        raise ValueError(f"Project {project_id} not found")

    # ---- Step 1: Clean draft ----
    finalized = await clean_scope(db, scope_data, project=project)
    overview = finalized.get("overview", {})

    # ---- Step 2: Update project metadata ----
    if overview:
        project.name = overview.get("Project Name") or project.name
        project.domain = overview.get("Domain") or project.domain
        project.complexity = overview.get("Complexity") or project.complexity
        project.tech_stack = overview.get("Tech Stack") or project.tech_stack
        project.use_cases = overview.get("Use Cases") or project.use_cases
        project.compliance = overview.get("Compliance") or project.compliance
        project.duration = str(overview.get("Duration") or project.duration)
        await db.commit()
        await db.refresh(project)

    # ---- Step 3: Save finalized_scope.json ----
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project_id,
            models.ProjectFile.file_name == "finalized_scope.json"
        )
    )
    old_file = result.scalars().first()
    if old_file:
        logger.info(f" Overwriting existing finalized_scope.json for project {project_id}")
    else:
        old_file = models.ProjectFile(
            project_id=project_id,
            file_name="finalized_scope.json",
        )

    blob_name = f"{PROJECTS_BASE}/{project_id}/finalized_scope.json"
    await azure_blob.upload_bytes(
        json.dumps(finalized, ensure_ascii=False, indent=2).encode("utf-8"),
        blob_name,
        overwrite=True,
    )

    old_file.file_path = blob_name
    db.add(old_file)
    await db.commit()
    await db.refresh(old_file)

    logger.info(f" Finalized scope saved (no LLM) for project {project_id}")
    return old_file, {**finalized, "_finalized": True}
